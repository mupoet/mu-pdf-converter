#!/usr/bin/env python3
"""
pdf_to_pptx.py - 将 PDF 高保真转换为可编辑 PowerPoint
用法：python pdf_to_pptx.py input.pdf [--outfile output.pptx] [--slide-size 16:9|4:3|A4]

每页 PDF 处理策略（分层叠加）：
  Layer 1: 位图（嵌入图片，保持原格式 PNG/JPEG）
  Layer 2: 矢量图（SVG，Office 2019+ 可编辑）
  Layer 3: 表格（原生可编辑 Table 对象）
  Layer 4: 文本框（可编辑，保留字体/颜色/加粗/斜体）

坐标系：
  - fitz/pdfplumber: pt，原点左上角
  - pptx EMU: 1pt = 12700 EMU
"""

import argparse
import io
import sys
import tempfile
from pathlib import Path

# ── 依赖检查 ──────────────────────────────────
_missing = []
try:
    import fitz  # PyMuPDF
except ImportError:
    _missing.append("pymupdf")
try:
    import pdfplumber
except ImportError:
    _missing.append("pdfplumber")
try:
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError:
    _missing.append("python-pptx")

if _missing:
    print(
        f"[错误] 缺少依赖：{', '.join(_missing)}\n"
        f"请运行：pip install {' '.join(_missing)}",
        file=sys.stderr,
    )
    sys.exit(1)

sys.path.insert(0, str(Path(__file__).parent))
from utils import (  # noqa: E402
    generate_output_path,
    is_in_table_region,
    bbox_overlap,
    fitz_color_to_rgb,
    map_font_name,
    is_scanned_pdf,
    check_scanned_and_warn,
    PT_TO_EMU,
)


# ─────────────────────────────────────────────
# Slide 尺寸预设
# ─────────────────────────────────────────────

SLIDE_SIZES = {
    "16:9": (Emu(9144000), Emu(5143500)),   # 25.4cm × 14.29cm（标准宽屏）
    "4:3":  (Emu(9144000), Emu(6858000)),   # 25.4cm × 19.05cm（标准）
    "A4":   (Emu(9906000), Emu(7772160)),   # 27.52cm × 21.59cm（A4 横向）
    "A4v":  (Emu(6858000), Emu(9906000)),   # A4 竖向
}

DEFAULT_SLIDE_SIZE = "pdf"  # 默认跟 PDF 页面尺寸完全一致


# ─────────────────────────────────────────────
# 核心转换函数
# ─────────────────────────────────────────────

def pdf_to_pptx(
    pdf_path: str,
    outfile: str = None,
    slide_size: str = DEFAULT_SLIDE_SIZE,
    auto_translate: bool = True,
    verbose_translate: bool = False,
) -> list[str]:
    """
    将 PDF 每页转换为 PowerPoint 幻灯片。

    参数:
        pdf_path: 输入 PDF 路径
        outfile: 输出 pptx 路径（可选）
        slide_size: 幻灯片尺寸预设 'pdf'（默认，跟PDF一致）|'16:9'|'4:3'|'A4'|'A4v'
        auto_translate: 是否自动检测外文并生成中文版（默认 True）
        verbose_translate: 打印逐条翻译日志
    返回:
        输出文件路径列表（外文 PDF 时返回 [原版路径, 中文版路径]，否则返回 [路径]）
    """
    from translate_utils import is_foreign_pdf, batch_translate

    pdf_path = str(Path(pdf_path).resolve())
    if not Path(pdf_path).exists():
        print(f"[错误] 文件不存在：{pdf_path}", file=sys.stderr)
        sys.exit(1)

    output_path = generate_output_path(pdf_path, ".pptx", outfile=outfile)

    # 检测扫描件
    scanned = is_scanned_pdf(pdf_path)
    if scanned:
        print("[警告] 检测到扫描件，将仅插入位图（无文本层）", file=sys.stderr)

    fitz_doc = fitz.open(pdf_path)
    total_pages = len(fitz_doc)

    # 确定 Slide 尺寸
    if slide_size == "pdf" or slide_size not in SLIDE_SIZES:
        first_page = fitz_doc[0]
        pdf_w0 = first_page.rect.width
        pdf_h0 = first_page.rect.height
        slide_w = Emu(int(pdf_w0 * PT_TO_EMU))
        slide_h = Emu(int(pdf_h0 * PT_TO_EMU))
        size_label = f"PDF原始({pdf_w0:.0f}×{pdf_h0:.0f}pt)"
    else:
        slide_w, slide_h = SLIDE_SIZES[slide_size]
        size_label = slide_size

    # ── 语言检测（提取全文本用于判断）──────────────
    need_translate = False
    if auto_translate and not scanned:
        all_texts = []
        for page_idx in range(total_pages):
            page_text = fitz_doc[page_idx].get_text("text")
            all_texts.append(page_text)
        need_translate = is_foreign_pdf(all_texts)
        if need_translate:
            print(f"[翻译] 检测到外文 PDF（CJK 占比 < 5%），将同时生成中文版", file=sys.stderr)
        else:
            print(f"[翻译] 非全外文 PDF，跳过翻译", file=sys.stderr)

    # ── 预提取每页文本块（供两版本复用）──────────
    # 结构: all_page_blocks[page_idx] = list of (block_data, text_str)
    all_page_blocks: list = []
    if not scanned:
        for page_idx in range(total_pages):
            fitz_page = fitz_doc[page_idx]
            page_dict = fitz_page.get_text("dict", sort=True)
            blocks = page_dict.get("blocks", [])
            page_texts = []
            for block in blocks:
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        t = span.get("text", "").strip()
                        if t:
                            page_texts.append(t)
            all_page_blocks.append(page_texts)
    else:
        all_page_blocks = [[] for _ in range(total_pages)]

    # ── 批量翻译（一次性，避免多次网络调用）────────
    MAX_TRANSLATE_BLOCKS = 5000  # 翻译文本块上限，防止超大PDF耗尽API配额

    translated_map: dict[int, dict[str, str]] = {}  # page_idx → {原文: 译文}
    if need_translate:
        all_texts_flat = []
        page_text_index = []  # (page_idx, text)
        for page_idx, page_texts in enumerate(all_page_blocks):
            for t in page_texts:
                all_texts_flat.append(t)
                page_text_index.append((page_idx, t))

        if len(all_texts_flat) > MAX_TRANSLATE_BLOCKS:
            print(f"[翻译] 文本块数 {len(all_texts_flat)} 超过上限 {MAX_TRANSLATE_BLOCKS}，截断处理", file=sys.stderr)
            all_texts_flat = all_texts_flat[:MAX_TRANSLATE_BLOCKS]
            page_text_index = page_text_index[:MAX_TRANSLATE_BLOCKS]

        print(f"[翻译] 共 {len(all_texts_flat)} 个文本块，批量翻译中...", file=sys.stderr)
        translated_flat = batch_translate(
            all_texts_flat,
            verbose=verbose_translate,
        )

        for i, (page_idx, orig) in enumerate(page_text_index):
            if page_idx not in translated_map:
                translated_map[page_idx] = {}
            translated_map[page_idx][orig] = translated_flat[i]
        print(f"[翻译] 完成", file=sys.stderr)

    # ── 生成 PPTX（原版，以及可选中文版）────────
    output_paths = []
    versions = [("原版", None)]
    if need_translate:
        stem = Path(output_path).stem
        zh_path = str(Path(output_path).parent / f"{stem}_中文版.pptx")
        versions.append(("中文版", zh_path))

    for version_label, version_outpath in versions:
        this_outpath = version_outpath if version_outpath else output_path
        translate_map_this = translated_map if (version_label == "中文版") else {}

        prs = Presentation()
        prs.slide_width = slide_w
        prs.slide_height = slide_h
        print(f"\n[信息] 生成{version_label}，共 {total_pages} 页，Slide 尺寸：{size_label}")

        with pdfplumber.open(pdf_path) as plumber_pdf:
            for page_idx in range(total_pages):
                fitz_page = fitz_doc[page_idx]
                plumber_page = plumber_pdf.pages[page_idx]
                page_num = page_idx + 1

                pdf_rect = fitz_page.rect
                pdf_w_pt = pdf_rect.width
                pdf_h_pt = pdf_rect.height

                blank_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_layout)

                scale_x = int(slide_w) / (pdf_w_pt * PT_TO_EMU)
                scale_y = int(slide_h) / (pdf_h_pt * PT_TO_EMU)

                print(f"  [第 {page_num} 页] PDF 尺寸: {pdf_w_pt:.0f}×{pdf_h_pt:.0f}pt", end="")

                n_images = _add_embedded_images(
                    slide, fitz_page, fitz_doc, pdf_w_pt, pdf_h_pt, slide_w, slide_h, scale_x, scale_y
                )
                n_svg = _add_svg_layer(
                    slide, fitz_page, pdf_w_pt, pdf_h_pt, slide_w, slide_h
                )

                table_bboxes = []
                n_tables = 0
                if not scanned:
                    page_translate = translate_map_this.get(page_idx, {})
                    table_bboxes, n_tables = _add_tables(
                        slide, plumber_page, pdf_w_pt, pdf_h_pt, slide_w, slide_h,
                        scale_x, scale_y, translate_map=page_translate
                    )

                n_texts = 0
                if not scanned:
                    page_translate = translate_map_this.get(page_idx, {})
                    n_texts = _add_text_boxes(
                        slide, fitz_page, table_bboxes, pdf_w_pt, pdf_h_pt,
                        slide_w, slide_h, scale_x, scale_y, translate_map=page_translate
                    )

                plumber_page.close()  # 修复内存泄漏
                print(f" | 图片={n_images} SVG={n_svg} 表格={n_tables} 文本框={n_texts}")

        prs.save(this_outpath)
        print(f"[完成] {version_label} → {this_outpath}")
        output_paths.append(this_outpath)

    fitz_doc.close()
    return output_paths


# ─────────────────────────────────────────────
# Layer 1: 嵌入位图
# ─────────────────────────────────────────────

def _add_embedded_images(slide, fitz_page, fitz_doc, pdf_w, pdf_h, slide_w, slide_h,
                          scale_x, scale_y) -> int:
    """提取 PDF 中嵌入的原始位图并插入 slide"""
    count = 0
    try:
        image_list = fitz_page.get_images(full=True)
    except Exception:
        return 0

    for img_info in image_list:
        xref = img_info[0]
        try:
            img_data = fitz_doc.extract_image(xref)
            if not img_data:
                continue

            img_bytes = img_data.get("image", b"")
            img_ext = img_data.get("ext", "png").lower()
            if not img_bytes:
                continue

            # 获取图片在页面上的位置（可能有多个）
            rects = fitz_page.get_image_rects(xref)
            if not rects:
                continue

            for rect in rects:
                left, top, width, height = _pdf_rect_to_slide_emu(
                    rect, pdf_w, pdf_h, slide_w, slide_h, scale_x, scale_y
                )
                if width <= 0 or height <= 0:
                    continue

                # python-pptx 不支持 JPEG2000/jp2，转码为 PNG 后再插入
                SUPPORTED_EXTS = {"png", "jpeg", "jpg", "bmp", "gif", "tiff", "wmf"}
                if img_ext not in SUPPORTED_EXTS:
                    try:
                        pix = fitz.Pixmap(fitz_doc, xref)
                        if pix.n > 4:  # CMYK → RGB
                            pix = fitz.Pixmap(fitz.csRGB, pix)
                        img_bytes = pix.tobytes("png")
                        img_ext = "png"
                    except Exception:
                        pass  # 转码失败则下面 add_picture 会抛出

                img_stream = io.BytesIO(img_bytes)
                try:
                    slide.shapes.add_picture(img_stream, left, top, width, height)
                    count += 1
                except Exception as e:
                    print(f"\n    [警告] 图片插入失败 (xref={xref}): {e}", file=sys.stderr)

        except Exception as e:
            print(f"\n    [警告] 图片提取失败 (xref={xref}): {e}", file=sys.stderr)
            continue

    return count


# ─────────────────────────────────────────────
# Layer 2: SVG 矢量图（整页铺底，可编辑）
# ─────────────────────────────────────────────

def _add_svg_layer(slide, fitz_page, pdf_w, pdf_h, slide_w, slide_h) -> int:
    """
    检测矢量路径，若存在则导出整页 SVG 并插入 slide。

    注意：python-pptx 插入 SVG 需要直接操作 XML，
    因为标准 API 不支持 SVG。Office 2019+ 支持 SVG 编辑。
    """
    try:
        drawings = fitz_page.get_drawings()
    except Exception:
        return 0

    if not drawings:
        return 0

    # 计算所有矢量元素的合并 bbox
    all_rects = []
    for d in drawings:
        rect = d.get("rect")
        if rect:
            all_rects.append(rect)

    if not all_rects:
        return 0

    # 导出整页 SVG
    try:
        svg_str = fitz_page.get_svg_image(matrix=fitz.Identity)
        svg_bytes = _strip_svg_text(svg_str.encode("utf-8"))
    except Exception as e:
        print(f"\n    [警告] SVG 导出失败: {e}", file=sys.stderr)
        return 0

    # 整页覆盖放置（整页 SVG 尺寸 = slide 尺寸）
    _insert_svg_to_slide(slide, svg_bytes, 0, 0, int(slide_w), int(slide_h))
    return 1


def _strip_svg_text(svg_bytes: bytes) -> bytes:
    """
    从 SVG 中删除所有文字节点（<text>、<tspan>、<flowRoot> 等），
    只保留图形/路径/矩形，避免与 pptx 文本框层重叠。
    """
    try:
        from lxml import etree
        NS = "http://www.w3.org/2000/svg"
        tree = etree.fromstring(svg_bytes)
        # 删除所有文字相关元素
        for tag in ("text", "tspan", "textPath", "flowRoot", "flowPara",
                    "flowSpan", "altGlyph", "altGlyphDef"):
            for elem in tree.iter(f"{{{NS}}}{tag}"):
                parent = elem.getparent()
                if parent is not None:
                    parent.remove(elem)
        return etree.tostring(tree, encoding="utf-8", xml_declaration=True)
    except Exception:
        # 解析失败原样返回
        return svg_bytes


def _insert_svg_to_slide(slide, svg_bytes: bytes, left: int, top: int,
                          width: int, height: int):
    """
    将 SVG 作为可编辑矢量图插入 pptx slide。

    实现方式：
    - 用 pptx.opc.package.Part 创建 SVG Part（content-type: image/svg+xml）
    - 通过 slide.part.relate_to() 注册关系，获取 rId
    - 手动构建 p:pic XML 元素并追加到 slide._spTree
    - Office 2019+ / LibreOffice 可直接编辑该 SVG 层

    注意：每页 SVG 的 partname 必须唯一，用 id(svg_bytes) 作区分。
    """
    from pptx.opc.package import Part
    from pptx.opc.packuri import PackURI
    from lxml import etree

    IMG_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
    SVG_CT = "image/svg+xml"

    sp = slide.part
    pkg = sp._package

    # 注册 SVG Part（partname 必须唯一）
    partname = PackURI(f"/ppt/media/svg_{abs(hash(svg_bytes)) % 10**9}.svg")
    svg_part = Part(partname, SVG_CT, pkg, svg_bytes)
    rId = sp.relate_to(svg_part, IMG_REL)

    pic_xml = (
        '<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        f' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        '<p:nvPicPr>'
        f'<p:cNvPr id="{abs(hash(svg_bytes)) % 9000 + 100}" name="SVGLayer"/>'
        '<p:cNvPicPr/>'
        '<p:nvPr/>'
        '</p:nvPicPr>'
        '<p:blipFill>'
        f'<a:blip r:embed="{rId}"/>'
        '<a:stretch><a:fillRect/></a:stretch>'
        '</p:blipFill>'
        '<p:spPr>'
        '<a:xfrm>'
        f'<a:off x="{left}" y="{top}"/>'
        f'<a:ext cx="{width}" cy="{height}"/>'
        '</a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        '</p:spPr>'
        '</p:pic>'
    )
    pic_elem = etree.fromstring(pic_xml)
    slide.shapes._spTree.append(pic_elem)


# ─────────────────────────────────────────────
# Layer 3: 原生可编辑表格
# ─────────────────────────────────────────────

def _add_tables(slide, plumber_page, pdf_w, pdf_h, slide_w, slide_h,
                scale_x, scale_y, translate_map: dict = None) -> tuple:
    """
    用 pdfplumber 识别表格并插入原生 pptx 表格。
    translate_map: {原文: 译文}，非空时替换单元格文本为译文
    返回 (table_bboxes, count)
    """
    table_bboxes = []
    count = 0

    # 严格表格识别参数——减少把文本段落误识别为表格
    TABLE_SETTINGS = {
        "vertical_strategy": "lines",       # 只用实线检测列（不用文字对齐推断）
        "horizontal_strategy": "lines",     # 只用实线检测行
        "snap_tolerance": 3,                # 线段对齐容差（默认3，不放宽）
        "join_tolerance": 3,
        "edge_min_length": 20,              # 线段最短 20pt 才算边框
        "min_words_vertical": 3,            # 竖向至少 3 个词才认为是分隔线
        "min_words_horizontal": 1,
    }
    try:
        tables_meta = plumber_page.find_tables(table_settings=TABLE_SETTINGS)
        tables_data = plumber_page.extract_tables(table_settings=TABLE_SETTINGS)
    except Exception:
        return [], 0

    for t_meta, t_data in zip(tables_meta, tables_data):
        if not t_data:
            continue

        # 表格 bbox（pdfplumber 坐标：x0, top, x1, bottom）
        bbox = t_meta.bbox
        table_bboxes.append(bbox)

        rows = len(t_data)
        cols = max((len(r) for r in t_data), default=1)
        if rows == 0 or cols == 0:
            continue

        # ── 过滤规则 1：至少 3行×3列才算真表格 ──
        if rows < 3 or cols < 3:
            continue

        # ── 过滤规则 2：单元格最小尺寸（避免极细线误识别） ──
        bbox_w = bbox[2] - bbox[0]
        bbox_h = bbox[3] - bbox[1]
        cell_w = bbox_w / cols
        cell_h = bbox_h / rows
        if cell_w < 15 or cell_h < 8:
            continue

        # ── 过滤规则 3：内容密度 —— 超过70%单元格为空则不是真表格 ──
        all_cells = [str(c or "").strip() for row in t_data for c in row]
        empty_ratio = sum(1 for c in all_cells if not c) / max(len(all_cells), 1)
        if empty_ratio > 0.7:
            continue

        # 转换到 slide 坐标
        left, top, width, height = _pdf_bbox_to_slide_emu(
            bbox, pdf_w, pdf_h, slide_w, slide_h, scale_x, scale_y
        )
        if width <= 0 or height <= 0:
            continue

        try:
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            tbl = table_shape.table

            for r_idx, row_data in enumerate(t_data):
                for c_idx in range(cols):
                    cell_val = row_data[c_idx] if c_idx < len(row_data) else ""
                    cell_text = str(cell_val or "").replace("\n", " ").strip()
                    cell = tbl.cell(r_idx, c_idx)
                    # 翻译模式：用译文替换单元格文本
                    if translate_map and cell_text in translate_map:
                        cell_text = translate_map[cell_text]
                    cell.text = cell_text
                    # 首行加粗
                    if r_idx == 0:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.bold = True
                    # 单元格底纹无填充、线条无色（避免与 SVG 层重叠）
                    _clear_cell_fill_and_border(cell)
            count += 1
        except Exception as e:
            print(f"\n    [警告] 表格插入失败: {e}", file=sys.stderr)

    return table_bboxes, count


def _clear_cell_fill_and_border(cell):
    """
    将表格单元格设为：底纹无填充（透明）+ 边框无色。
    直接操作底层 XML，避免与 SVG 背景层重叠产生遮挡。
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    tc = cell._tc  # <a:tc> 元素

    # ── 1. 底纹无填充：设置 <a:tcPr> → <a:noFill> ──
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn("a:tcPr"))

    # 删除已有填充节点，再插 noFill
    for fill_tag in ("a:solidFill", "a:gradFill", "a:pattFill",
                     "a:blipFill", "a:grpFill", "a:noFill"):
        for old in tcPr.findall(qn(fill_tag)):
            tcPr.remove(old)
    etree.SubElement(tcPr, qn("a:noFill"))

    # ── 2. 四条边框无色：lnL/lnR/lnT/lnB ──
    NO_LINE_XML = (
        '<a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:noFill/>'
        '</a:ln>'
    )
    for border_attr in ("lnL", "lnR", "lnT", "lnB"):
        for old in tcPr.findall(qn(f"a:{border_attr}")):
            tcPr.remove(old)
        ln_elem = etree.fromstring(NO_LINE_XML)
        ln_elem.tag = qn(f"a:{border_attr}")
        tcPr.append(ln_elem)


# ─────────────────────────────────────────────
# Layer 4: 文本框
# ─────────────────────────────────────────────

def _add_text_boxes(slide, fitz_page, table_bboxes, pdf_w, pdf_h,
                    slide_w, slide_h, scale_x, scale_y,
                    translate_map: dict = None) -> int:
    """
    提取 PDF 文本块并插入可编辑文本框。
    跳过与表格区域重叠的文本块。
    translate_map: {原文: 译文}，非空时替换 span 文本为译文
    """
    count = 0
    try:
        text_dict = fitz_page.get_text("dict", sort=True)
    except Exception:
        return 0

    blocks = text_dict.get("blocks", [])

    for block in blocks:
        if block.get("type") != 0:  # 只处理文本块
            continue

        block_bbox = block.get("bbox")
        if not block_bbox:
            continue

        # 跳过表格区域
        if is_in_table_region(block_bbox, table_bboxes, threshold=0.3):
            continue

        # 转换坐标
        left, top, width, height = _pdf_bbox_to_slide_emu(
            block_bbox, pdf_w, pdf_h, slide_w, slide_h, scale_x, scale_y
        )
        # 最小尺寸保护
        width = max(width, int(Pt(10) * PT_TO_EMU))
        height = max(height, int(Pt(8) * PT_TO_EMU))

        # 添加文本框
        try:
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True

            first_para = True
            for line in block.get("lines", []):
                if first_para:
                    para = tf.paragraphs[0]
                    first_para = False
                else:
                    para = tf.add_paragraph()

                for span in line.get("spans", []):
                    span_text = span.get("text", "").rstrip("\n")
                    if not span_text:
                        continue

                    run = para.add_run()
                    # 翻译模式：用译文替换 span 文本
                    display_text = span_text.strip()
                    if translate_map and display_text in translate_map:
                        run.text = translate_map[display_text]
                    else:
                        run.text = span_text

                    # 字体属性
                    font_size = span.get("size", 11.0)
                    flags = span.get("flags", 0)
                    is_bold = bool(flags & 0b10000)   # bit 4: bold
                    is_italic = bool(flags & 0b000001) # bit 0: italic
                    color = fitz_color_to_rgb(span.get("color", 0))
                    font_name = map_font_name(span.get("font", ""))

                    run.font.size = Pt(font_size)
                    run.font.bold = is_bold
                    run.font.italic = is_italic
                    run.font.name = font_name
                    try:
                        run.font.color.rgb = RGBColor(*color)
                    except Exception:
                        pass

            count += 1
        except Exception as e:
            print(f"\n    [警告] 文本框插入失败: {e}", file=sys.stderr)

    return count


# ─────────────────────────────────────────────
# 坐标转换辅助
# ─────────────────────────────────────────────

def _pdf_rect_to_slide_emu(rect, pdf_w, pdf_h, slide_w, slide_h,
                            scale_x, scale_y) -> tuple:
    """将 fitz.Rect 转换为 slide EMU 坐标 (left, top, width, height)"""
    x0, y0, x1, y1 = rect.x0, rect.y0, rect.x1, rect.y1
    left  = int(x0 * PT_TO_EMU * scale_x)
    top   = int(y0 * PT_TO_EMU * scale_y)
    width = int((x1 - x0) * PT_TO_EMU * scale_x)
    height= int((y1 - y0) * PT_TO_EMU * scale_y)
    return left, top, width, height


def _pdf_bbox_to_slide_emu(bbox, pdf_w, pdf_h, slide_w, slide_h,
                            scale_x, scale_y) -> tuple:
    """将 (x0, y0, x1, y1) bbox 转换为 slide EMU 坐标"""
    x0, y0, x1, y1 = bbox
    left  = int(x0 * PT_TO_EMU * scale_x)
    top   = int(y0 * PT_TO_EMU * scale_y)
    width = int((x1 - x0) * PT_TO_EMU * scale_x)
    height= int((y1 - y0) * PT_TO_EMU * scale_y)
    return left, top, width, height


# ─────────────────────────────────────────────
# CLI 入口
# ─────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="将 PDF 高保真转换为可编辑 PowerPoint（.pptx）",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  python pdf_to_pptx.py report.pdf
  python pdf_to_pptx.py report.pdf --outfile output/result.pptx
  python pdf_to_pptx.py report.pdf --slide-size 4:3
  python pdf_to_pptx.py report.pdf --slide-size A4 --outfile report_A4.pptx

幻灯片尺寸:
  16:9  宽屏（默认，25.4×14.3cm）
  4:3   标准（25.4×19.1cm）
  A4    A4横向（27.5×21.6cm）

分层转换策略（由下到上叠加）:
  Layer 1: 嵌入位图 - 保持原始 PNG/JPEG 格式，无损质量
  Layer 2: 矢量图 SVG - 整页 SVG，Office 2019+ 可编辑路径
  Layer 3: 原生表格 - 可编辑的 Table 对象
  Layer 4: 文本框 - 可编辑，保留字体/颜色/加粗/斜体

注意:
  - 扫描件 PDF 只会包含 Layer 1（位图）
  - SVG 支持需要 Office 2019 或 LibreOffice 7+
  - 复杂多栏布局可能出现文字重叠
        """,
    )
    parser.add_argument("input", help="输入 PDF 文件路径")
    parser.add_argument(
        "--outfile", default=None,
        help="输出 pptx 文件路径（默认与 PDF 同名同目录）"
    )
    parser.add_argument(
        "--slide-size",
        choices=["pdf"] + list(SLIDE_SIZES.keys()),
        default=DEFAULT_SLIDE_SIZE,
        help="幻灯片尺寸：pdf（默认，跟PDF页面完全一致）|16:9|4:3|A4|A4v"
    )
    parser.add_argument(
        "--no-translate",
        action="store_true",
        default=False,
        help="禁用自动翻译检测（不生成中文版）"
    )
    parser.add_argument(
        "--verbose-translate",
        action="store_true",
        default=False,
        help="打印逐条翻译日志"
    )

    args = parser.parse_args()
    check_scanned_and_warn(args.input)
    output_paths = pdf_to_pptx(
        pdf_path=args.input,
        outfile=args.outfile,
        slide_size=args.slide_size,
        auto_translate=not args.no_translate,
        verbose_translate=args.verbose_translate,
    )
    if len(output_paths) > 1:
        print(f"\n已生成 {len(output_paths)} 个文件：")
        for p in output_paths:
            print(f"  → {p}")


if __name__ == "__main__":
    main()
